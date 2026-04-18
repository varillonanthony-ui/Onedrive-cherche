"""
OneDrive indexer — crawls all files and extracts text content for full-text search.
Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, MD.
"""
import io
import os
import json
import sqlite3
import logging
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from onedrive_client import OneDriveClient
from db import Database

logger = logging.getLogger(__name__)

# Extensions we attempt to extract text from
TEXT_EXTENSIONS = {
    "pdf", "docx", "doc", "xlsx", "xls", "csv",
    "pptx", "ppt", "txt", "md", "json", "xml", "html",
}

DELTA_LINK_FILE = ".delta_link.json"


def extract_text(file_bytes: bytes, extension: str) -> str:
    """Extract plain text from file bytes. Returns empty string on failure."""
    ext = extension.lower().lstrip(".")
    try:
        if ext == "pdf":
            return _extract_pdf(file_bytes)
        elif ext in ("docx", "doc"):
            return _extract_docx(file_bytes)
        elif ext in ("xlsx", "xls"):
            return _extract_xlsx(file_bytes)
        elif ext == "csv":
            return file_bytes.decode("utf-8", errors="replace")
        elif ext in ("pptx", "ppt"):
            return _extract_pptx(file_bytes)
        elif ext in ("txt", "md", "json", "xml", "html"):
            return file_bytes.decode("utf-8", errors="replace")
    except Exception as e:
        logger.debug(f"Extraction failed for .{ext}: {e}")
    return ""


def _extract_pdf(data: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except ImportError:
        pass
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        pass
    return ""


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


class OneDriveIndexer:
    def __init__(
        self,
        client: OneDriveClient,
        db_path: str = "onedrive_index.db",
        index_content: bool = True,
        max_file_size_mb: int = 20,
        num_workers: int = 4,
    ):
        self.client = client
        self.db = Database(db_path)
        self.index_content = index_content
        self.max_bytes = max_file_size_mb * 1024 * 1024
        self.num_workers = num_workers

    def _load_delta_link(self) -> str | None:
        if os.path.exists(DELTA_LINK_FILE):
            with open(DELTA_LINK_FILE) as f:
                return json.load(f).get("link")
        return None

    def _save_delta_link(self, link: str):
        with open(DELTA_LINK_FILE, "w") as f:
            json.dump({"link": link, "saved_at": datetime.utcnow().isoformat()}, f)

    def run(self, start_path: str = "/", delta: bool = False, callback=None) -> dict:
        """
        Full index or delta update.
        callback(current, total, filename, error=False) called for progress.
        Returns stats dict.
        """
        self.db.init()

        delta_link = self._load_delta_link() if delta else None
        items, new_delta = self.client.list_items(start_path, delta_link=delta_link)

        if new_delta:
            self._save_delta_link(new_delta)

        # Only files (not folders)
        files = [i for i in items if "file" in i]
        deleted = [i for i in items if i.get("deleted")]

        # Remove deleted items
        for d in deleted:
            self.db.delete_item(d.get("id", ""))

        total = len(files)
        indexed = 0
        errors = 0

        def process_file(item: dict) -> tuple[bool, str]:
            try:
                item_id = item["id"]
                name = item.get("name", "")
                path = self.client.get_item_path(item)
                ext = Path(name).suffix.lower().lstrip(".")
                size = item.get("size", 0)
                modified = item.get("lastModifiedDateTime", "")

                content_text = ""
                if self.index_content and ext in TEXT_EXTENSIONS and size <= self.max_bytes:
                    file_bytes = self.client.download_file(item_id, max_bytes=self.max_bytes)
                    if file_bytes:
                        content_text = extract_text(file_bytes, ext)

                self.db.upsert_file(
                    item_id=item_id,
                    name=name,
                    path=path,
                    extension=ext,
                    size_bytes=size,
                    modified_at=modified,
                    content_text=content_text,
                )
                return True, name
            except Exception as e:
                logger.warning(f"Error processing {item.get('name')}: {e}")
                return False, item.get("name", "?")

        with ThreadPoolExecutor(max_workers=self.num_workers) as pool:
            futures = {pool.submit(process_file, f): f for f in files}
            done = 0
            for future in as_completed(futures):
                done += 1
                ok, fname = future.result()
                if ok:
                    indexed += 1
                else:
                    errors += 1
                if callback:
                    callback(done, total, fname, error=not ok)

        return {
            "total": total,
            "indexed": indexed,
            "errors": errors,
            "deleted": len(deleted),
            "timestamp": datetime.utcnow().isoformat(),
        }
